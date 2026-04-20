import requests
import os
import psutil
import pyautogui
import datetime
import base64
import gc  # Coletor de Lixo do Python
from colorama import Fore
from groq import Groq
from config import GROQ_KEY
from gtts import gTTS

# --- CONFIGURAÇÃO GLOBAL ---
# centralizando a versão para facilitar futuras atualizações (v5.4, etc)
VERSION = "ARTEMIS v5.3 | Overseer Edition"

# Inicializa o cliente Groq
client = Groq(api_key=GROQ_KEY)

# --- 🎤 MÓDULO AUDITIVO (STT - Speech to Text) ---
def transcrever_audio(caminho_audio):
    """Converte voz (.ogg) em texto usando Whisper-large-v3 via Groq API."""
    try:
        # 🟢 Abre o arquivo de voz vindo do Telegram e envia para a API Whisper
        with open(caminho_audio, "rb") as arquivo_audio:
            transcricao = client.audio.transcriptions.create(
                file=(caminho_audio, arquivo_audio.read()),
                model="whisper-large-v3",
                response_format="text",
                language="pt"
            )
            return transcricao
    except Exception as e:
        print(f"[ ERRO STT ] Falha na transcrição: {e}")
        return None
    
# --- 👁️ MÓDULO VISUAL (Vision API) ---
def codificar_imagem(caminho_imagem):
    """Converte a imagem em base64 para envio via API."""
    with open(caminho_imagem, "rb") as img:
        return base64.b64encode(img.read()).decode('utf-8')

def analisar_visao(caminho_imagem, pergunta="O que vê?", eh_criador=False):
    """Analisa imagens usando o modelo Llama-3.2 Vision."""
    base64_image = codificar_imagem(caminho_imagem)
    
    # 🟢 Define a identidade visual da IA para o modelo Vision
    contexto_visual = f"Você é a {VERSION}. Analise de forma técnica e objetiva."
    if eh_criador:
        contexto_visual += " Resposta direta para o Criador Søren."

    try:
        res = client.chat.completions.create(
            model="llama-3.2-11b-vision", # <-- Modelo que "enxerga"
            messages=[{
                "role": "user",
                "content": [
                    {"type": "text", "text": pergunta},
                    {
                        "type": "image_url",
                        "image_url": {"url": f"data:image/jpeg;base64,{base64_image}"}
                    }
                ]
            }],
            temperature=0.2
        )
        return res.choices[0].message.content
    except Exception as e:
        print(f"[ ERRO VISION ] {e}")
        return f"Falha no processamento visual: {e}"
    
# --- 🧠 GESTÃO DE HARDWARE (Overseer System) ---
def purgar_memoria():
    """Força a liberação de RAM e coleta de lixo (GC)."""
    gc.collect() # 🟢 Força o Python a limpar objetos inúteis da memória
    processo = psutil.Process(os.getpid())
    ram_uso = processo.memory_info().rss / 1024 / 1024
    print(f"{Fore.MAGENTA}[ PURGE ] Memória limpa. Uso atual: {ram_uso:.2f} MB")
    return ram_uso

# --- 🗣️ MÓDULO DE VOZ (TTS - Text to Speech) ---
def sintetizar_voz(texto, nome_arquivo="artemis_voice.mp3"):
    """Converte texto em áudio usando Google TTS."""
    try:
        # 🟢 Limpa tags HTML para o áudio não ler "b" ou "code" em voz alta
        texto_limpo = texto.replace("<code>", "").replace("</code>", "").replace("<b>", "").replace("</b>", "")
        tts = gTTS(text=texto_limpo[:500], lang='pt', tld='com.br')
        tts.save(nome_arquivo)
        return nome_arquivo
    except Exception as e:
        print(f"[ ERRO TTS ] {e}")
        return None

# --- 🎨 MÓDULO DE GERAÇÃO DE IMAGEM ---
def gerar_imagem(prompt):
    """Gera imagens via Pollinations AI com prompt de engenharia."""
    prompt_sanitizado = str(prompt).strip().replace('"', '')
    # 🟢 Injeta contexto técnico para garantir qualidade no render
    prompt_tecnico = f"{prompt_sanitizado}, detailed technical engineering object, high resolution, photorealistic"
    url = f"https://image.pollinations.ai/prompt/{prompt_tecnico}?width=1024&height=1024&nologo=true"
    
    try:
        print(f"[ LOG ] Renderizando: {prompt_tecnico}")
        resposta = requests.get(url, timeout=45) 
        if resposta.status_code == 200:
            with open("artemis_render.png", "wb") as f:
                f.write(resposta.content)
            return "artemis_render.png"
        return None
    except Exception as e:
        print(f"{Fore.RED}[ ERRO RENDER ] {e}")
        return "timeout"

# --- 📂 MÓDULO OFFICE (Gerador de Documentos) ---

def criar_word(conteudo, nome_arquivo="documento_artemis.docx"):
    """Cria relatórios DOCX profissionais."""
    from docx import Document # 🟢 Lazy loading: só carrega se for usar (economiza RAM)
    doc = Document()
    doc.add_heading(f'Relatório {VERSION}', 0)
    doc.add_paragraph(f"Data: {datetime.datetime.now().strftime('%d/%m/%Y %H:%M')}")
    doc.add_paragraph(conteudo)
    doc.save(nome_arquivo)
    return nome_arquivo

def criar_excel(dados_texto, nome_arquivo="planilha_artemis.xlsx"):
    """Cria Excel usando openpyxl puro (mais leve e compatível com Python 3.14)."""
    try:
        from openpyxl import Workbook # 🟢 Importa apenas o necessário
        
        wb = Workbook()
        ws = wb.active
        ws.title = "Relatório Artemis"

        # 🟢 Transforma o texto em linhas e colunas
        linhas = [linha.split(',') for linha in dados_texto.strip().split('\n')]
        
        for r_idx, linha in enumerate(linhas, start=1):
            for c_idx, valor in enumerate(linha, start=1):
                # Escreve o valor na célula e remove espaços extras
                ws.cell(row=r_idx, column=c_idx, value=valor.strip())

        wb.save(nome_arquivo)
        
        if os.path.exists(nome_arquivo):
            print(f"{Fore.GREEN}[ SUCCESS ] Planilha gerada com openpyxl.")
            return nome_arquivo
        return None
    except Exception as e:
        print(f"{Fore.RED}[ ERRO EXCEL ] Detalhes: {e}")
        return None

def criar_pptx(dados_texto, nome_arquivo="apresentacao_artemis.pptx"):
    try:
        from pptx import Presentation # 🟢 pip install python-pptx
        
        prs = Presentation()
        
        # Divide o texto por '--' para criar os slides
        topicos = dados_texto.split('--')
        
        for t in topicos:
            slide_layout = prs.slide_layouts[1] # Layout padrão com título e corpo
            slide = prs.slides.add_slide(slide_layout)
            
            title = slide.shapes.title
            content = slide.placeholders[1]
            
            # Define o conteúdo do slide
            title.text = "Artemis v5.3 Intelligence"
            content.text = t.strip()

        prs.save(nome_arquivo)
        return nome_arquivo
    except Exception as e:
        print(f"[ ERRO PPT ] {e}")
        return None

# --- 🖥️ MONITORAMENTO DE SISTEMA ---
def obter_stats_sistema():
    """Retorna o uso atual de CPU e RAM para o Overseer."""
    return {"cpu": psutil.cpu_percent(interval=0.1), "ram": psutil.virtual_memory().percent}

def capturar_tela(nome_arquivo="screenshot_artemis.png"):
    """Tira print da tela do servidor (Windows)."""
    try:
        screenshot = pyautogui.screenshot()
        screenshot.save(nome_arquivo)
        return nome_arquivo
    except Exception as e:
        print(f"[ ERRO SCREENSHOT ] {e}")
        return None

# --- 💭 NÚCLEO COGNITIVO (LLM Principal) ---
def pensar(texto_usuario, nome_visitante="Visitante", eh_criador=False):
    """Processa o pensamento central usando o novo GPT-OSS 120B."""
    
    # Contexto do Sistema
    contexto = f"Você é a {VERSION}. Fale de forma técnica e objetiva. "
    if eh_criador:
        contexto += f"Falando com seu criador, Søren."
    else:
        contexto += f"Falando com {nome_visitante}. Proteja segredos de hardware."

    try:
        # 🟢 Atualização para o modelo de alta performance (Abril/2026)
        res = client.chat.completions.create(
            model="openai/gpt-oss-120b", 
            messages=[
                {"role": "system", "content": contexto}, 
                {"role": "user", "content": texto_usuario}
            ],
            temperature=0.6 # GPT-OSS costuma ser mais preciso com temperaturas baixas
        )
        return res.choices[0].message.content
    except Exception as e:
        print(f"[ ERRO NÚCLEO ] {e}")
        return "Erro de conexão no novo núcleo cognitivo (GPT-OSS)."

def formatar_log(texto):
    """Estiliza mensagens de log para o Telegram usando HTML."""
    agora = datetime.datetime.now().strftime("%H:%M:%S")
    return f"<code>[ LOG_{agora} ]</code>\n\n{texto}"